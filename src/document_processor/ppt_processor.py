from pptx import Presentation
from typing import List, Dict
import base64

class PPTProcessor:
    def extract_content(self, ppt_path: str) -> Dict:
        prs = Presentation(ppt_path)
        content = {
            'slides': [],
            'metadata': {
                'title': prs.core_properties.title or '',
                'author': prs.core_properties.author or '',
                'slide_count': len(prs.slides),
                'source': ppt_path
            }
        }
        
        for slide_num, slide in enumerate(prs.slides):
            slide_content = {
                'slide_number': slide_num + 1,
                'title': '',
                'text': '',
                'images': [],
                'notes': ''
            }
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    if shape.text.strip():
                        if not slide_content['title'] and len(shape.text.strip().split('\n')[0]) < 100:
                            slide_content['title'] = shape.text.strip().split('\n')[0]
                        slide_content['text'] += shape.text + '\n'
                
                # Extract images
                if shape.shape_type == 13:  # Picture
                    try:
                        image_data = shape.image.blob
                        slide_content['images'].append({
                            'data': base64.b64encode(image_data).decode(),
                            'format': shape.image.ext
                        })
                    except:
                        pass
            
            # Extract speaker notes
            if slide.has_notes_slide:
                slide_content['notes'] = slide.notes_slide.notes_text_frame.text
            
            content['slides'].append(slide_content)
        
        return content
    
    def create_chunks(self, ppt_content: Dict) -> List[Dict]:
        chunks = []
        
        for slide in ppt_content['slides']:
            # Create text chunk for each slide
            text_content = f"Slide {slide['slide_number']}: {slide['title']}\n{slide['text']}"
            if slide['notes']:
                text_content += f"\nNotes: {slide['notes']}"
            
            chunks.append({
                'text': text_content,
                'metadata': {
                    **ppt_content['metadata'],
                    'slide_number': slide['slide_number'],
                    'slide_title': slide['title'],
                    'chunk_type': 'slide',
                    'has_images': len(slide['images']) > 0
                }
            })
            
            # Create separate chunks for images with descriptions
            for img_idx, image in enumerate(slide['images']):
                chunks.append({
                    'text': f"Image from slide {slide['slide_number']}: {slide['title']}",
                    'image_data': image['data'],
                    'metadata': {
                        **ppt_content['metadata'],
                        'slide_number': slide['slide_number'],
                        'image_index': img_idx,
                        'chunk_type': 'image'
                    }
                })
        
        return chunks
