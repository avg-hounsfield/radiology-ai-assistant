#!/usr/bin/env python3
"""
Comprehensive Image Processing and Embedding System for ECHO
Handles image extraction, processing, embedding, and search capabilities
"""

import os
import json
import logging
import hashlib
from pathlib import Path
from datetime import datetime
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass, asdict
import base64
import io

# Core libraries
try:
    import cv2
    CV2_AVAILABLE = True
except ImportError:
    CV2_AVAILABLE = False
    cv2 = None

import numpy as np
from PIL import Image, ImageOps
import fitz  # PyMuPDF for PDF processing
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ML/AI libraries
try:
    import torch
    from transformers import CLIPProcessor, CLIPModel
    CLIP_AVAILABLE = True
except ImportError:
    CLIP_AVAILABLE = False
    logging.warning("CLIP not available - visual embeddings will be disabled")

@dataclass
class RadiologyImage:
    """Represents a radiology image with metadata and embeddings"""
    image_id: str
    file_path: str
    source_document: str = ""
    page_number: int = -1
    slide_number: int = -1

    # Image properties
    width: int = 0
    height: int = 0
    format: str = ""
    size_bytes: int = 0

    # Medical metadata
    modality: str = ""  # CT, MRI, X-ray, etc.
    body_part: str = ""  # chest, abdomen, brain, etc.
    pathology: List[str] = None
    tags: List[str] = None

    # Processing metadata
    created_date: str = ""
    last_modified: str = ""
    extracted_text: str = ""  # OCR or caption text

    # Embeddings
    visual_embedding: List[float] = None
    text_embedding: List[float] = None

    # Display
    thumbnail_base64: str = ""

    def __post_init__(self):
        if self.pathology is None:
            self.pathology = []
        if self.tags is None:
            self.tags = []
        if self.created_date == "":
            self.created_date = datetime.now().isoformat()

    def to_dict(self):
        return asdict(self)

    @classmethod
    def from_dict(cls, data: dict):
        return cls(**data)

class ImageProcessor:
    """Core image processing and analysis functionality"""

    def __init__(self, data_dir: str = "data/images"):
        self.data_dir = Path(data_dir)
        self.data_dir.mkdir(parents=True, exist_ok=True)

        self.thumbnails_dir = self.data_dir / "thumbnails"
        self.thumbnails_dir.mkdir(exist_ok=True)

        # Initialize CLIP model if available
        self.clip_model = None
        self.clip_processor = None
        if CLIP_AVAILABLE:
            try:
                self.clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
                self.clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
                logging.info("✅ CLIP model loaded for visual embeddings")
            except Exception as e:
                logging.warning(f"⚠️ Failed to load CLIP model: {e}")
                CLIP_AVAILABLE = False

    def process_image(self, image_path: str, source_doc: str = "", page_num: int = -1) -> Optional[RadiologyImage]:
        """Process a single image file and extract metadata"""
        try:
            # Load and analyze image
            image = Image.open(image_path)
            image_array = np.array(image)

            # Generate unique ID
            with open(image_path, 'rb') as f:
                image_hash = hashlib.md5(f.read()).hexdigest()

            # Extract basic properties
            file_stat = os.stat(image_path)

            # Create thumbnail and base64 encoding
            thumbnail = self._create_thumbnail(image)
            thumbnail_base64 = self._image_to_base64(thumbnail)

            # Detect modality and body part using filename and content analysis
            modality = self._detect_modality(image_path, image_array)
            body_part = self._detect_body_part(image_path, image_array)

            # Generate embeddings
            visual_embedding = self._generate_visual_embedding(image)

            # Create RadiologyImage object
            rad_image = RadiologyImage(
                image_id=image_hash,
                file_path=str(image_path),
                source_document=source_doc,
                page_number=page_num,
                width=image.width,
                height=image.height,
                format=image.format or "Unknown",
                size_bytes=file_stat.st_size,
                modality=modality,
                body_part=body_part,
                last_modified=datetime.fromtimestamp(file_stat.st_mtime).isoformat(),
                visual_embedding=visual_embedding,
                thumbnail_base64=thumbnail_base64
            )

            return rad_image

        except Exception as e:
            logging.error(f"Error processing image {image_path}: {e}")
            return None

    def _create_thumbnail(self, image: Image.Image, size: Tuple[int, int] = (150, 150)) -> Image.Image:
        """Create thumbnail with proper aspect ratio"""
        thumbnail = image.copy()
        thumbnail.thumbnail(size, Image.Resampling.LANCZOS)

        # Create a square thumbnail with padding
        thumb_with_padding = Image.new('RGB', size, (240, 240, 240))
        paste_x = (size[0] - thumbnail.width) // 2
        paste_y = (size[1] - thumbnail.height) // 2
        thumb_with_padding.paste(thumbnail, (paste_x, paste_y))

        return thumb_with_padding

    def _image_to_base64(self, image: Image.Image) -> str:
        """Convert PIL image to base64 string for web display"""
        buffer = io.BytesIO()
        image.save(buffer, format='JPEG', quality=85)
        img_str = base64.b64encode(buffer.getvalue()).decode()
        return f"data:image/jpeg;base64,{img_str}"

    def _detect_modality(self, file_path: str, image_array: np.ndarray) -> str:
        """Detect imaging modality from filename and image characteristics"""
        filename = os.path.basename(file_path).lower()

        # Keyword-based detection
        modality_keywords = {
            'CT': ['ct', 'computed', 'tomography', 'axial'],
            'MRI': ['mri', 'magnetic', 'resonance', 't1', 't2', 'flair', 'dwi'],
            'X-ray': ['xray', 'x-ray', 'radiograph', 'chest', 'ap', 'pa', 'lateral'],
            'Ultrasound': ['ultrasound', 'us', 'echo', 'doppler'],
            'Mammography': ['mammo', 'breast', 'cc', 'mlo'],
            'Nuclear': ['pet', 'spect', 'nuclear', 'bone scan', 'thyroid'],
            'Fluoroscopy': ['fluoro', 'barium', 'contrast']
        }

        for modality, keywords in modality_keywords.items():
            if any(keyword in filename for keyword in keywords):
                return modality

        # Image-based heuristics
        if len(image_array.shape) == 2 or (len(image_array.shape) == 3 and image_array.shape[2] == 1):
            # Grayscale - likely medical imaging
            mean_intensity = np.mean(image_array)
            if mean_intensity < 50:  # Very dark - might be nuclear medicine
                return "Nuclear"
            elif mean_intensity > 200:  # Very bright - might be X-ray
                return "X-ray"

        return "Unknown"

    def _detect_body_part(self, file_path: str, image_array: np.ndarray) -> str:
        """Detect anatomical region from filename"""
        filename = os.path.basename(file_path).lower()

        body_part_keywords = {
            'Chest': ['chest', 'thorax', 'lung', 'heart', 'mediastinum'],
            'Abdomen': ['abdomen', 'liver', 'kidney', 'pancreas', 'spleen'],
            'Pelvis': ['pelvis', 'bladder', 'prostate', 'uterus', 'ovary'],
            'Brain': ['brain', 'head', 'cranial', 'cerebral', 'skull'],
            'Spine': ['spine', 'vertebra', 'lumbar', 'cervical', 'thoracic'],
            'MSK': ['bone', 'joint', 'fracture', 'orthopedic', 'extremity'],
            'Breast': ['breast', 'mammo', 'cc', 'mlo']
        }

        for body_part, keywords in body_part_keywords.items():
            if any(keyword in filename for keyword in keywords):
                return body_part

        return "Unknown"

    def _generate_visual_embedding(self, image: Image.Image) -> Optional[List[float]]:
        """Generate visual embedding using CLIP model"""
        if not CLIP_AVAILABLE or not self.clip_model:
            return None

        try:
            # Preprocess image
            inputs = self.clip_processor(images=image, return_tensors="pt")

            # Generate embedding
            with torch.no_grad():
                image_features = self.clip_model.get_image_features(**inputs)
                embedding = image_features.squeeze().numpy().tolist()

            return embedding

        except Exception as e:
            logging.error(f"Error generating visual embedding: {e}")
            return None

class DocumentImageExtractor:
    """Extract images from PDFs and PowerPoint presentations"""

    def __init__(self, output_dir: str = "data/images/extracted"):
        self.output_dir = Path(output_dir)
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.processor = ImageProcessor()

    def extract_from_pdf(self, pdf_path: str) -> List[RadiologyImage]:
        """Extract all images from a PDF document"""
        extracted_images = []

        try:
            pdf_document = fitz.open(pdf_path)
            pdf_name = os.path.basename(pdf_path)

            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                image_list = page.get_images()

                for img_index, img in enumerate(image_list):
                    try:
                        # Extract image data
                        xref = img[0]
                        pix = fitz.Pixmap(pdf_document, xref)

                        if pix.n - pix.alpha < 4:  # Skip if not grayscale or RGB
                            # Save image
                            image_filename = f"{pdf_name}_page{page_num+1}_img{img_index+1}.png"
                            image_path = self.output_dir / image_filename

                            if pix.alpha:
                                pix = fitz.Pixmap(fitz.csRGB, pix)

                            pix.save(str(image_path))
                            pix = None

                            # Process extracted image
                            rad_image = self.processor.process_image(
                                str(image_path),
                                source_doc=pdf_path,
                                page_num=page_num + 1
                            )

                            if rad_image:
                                # Extract text from page for context
                                page_text = page.get_text()
                                rad_image.extracted_text = self._clean_text(page_text)
                                extracted_images.append(rad_image)

                    except Exception as e:
                        logging.error(f"Error extracting image {img_index} from page {page_num}: {e}")
                        continue

            pdf_document.close()
            logging.info(f"Extracted {len(extracted_images)} images from {pdf_path}")

        except Exception as e:
            logging.error(f"Error processing PDF {pdf_path}: {e}")

        return extracted_images

    def extract_from_ppt(self, ppt_path: str) -> List[RadiologyImage]:
        """Extract all images from a PowerPoint presentation"""
        extracted_images = []

        try:
            prs = Presentation(ppt_path)
            ppt_name = os.path.basename(ppt_path)

            for slide_num, slide in enumerate(prs.slides):
                for shape_index, shape in enumerate(slide.shapes):
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Extract image
                            image_stream = io.BytesIO(shape.image.blob)
                            image = Image.open(image_stream)

                            # Save image
                            image_filename = f"{ppt_name}_slide{slide_num+1}_img{shape_index+1}.png"
                            image_path = self.output_dir / image_filename
                            image.save(image_path)

                            # Process extracted image
                            rad_image = self.processor.process_image(
                                str(image_path),
                                source_doc=ppt_path,
                                page_num=slide_num + 1
                            )

                            if rad_image:
                                # Extract slide text for context
                                slide_text = self._extract_slide_text(slide)
                                rad_image.extracted_text = slide_text
                                rad_image.slide_number = slide_num + 1
                                extracted_images.append(rad_image)

                        except Exception as e:
                            logging.error(f"Error extracting image from slide {slide_num}: {e}")
                            continue

            logging.info(f"Extracted {len(extracted_images)} images from {ppt_path}")

        except Exception as e:
            logging.error(f"Error processing PowerPoint {ppt_path}: {e}")

        return extracted_images

    def _extract_slide_text(self, slide) -> str:
        """Extract text from a PowerPoint slide"""
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
        return self._clean_text(" ".join(text_runs))

    def _clean_text(self, text: str) -> str:
        """Clean and normalize extracted text"""
        if not text:
            return ""

        # Remove excessive whitespace
        text = " ".join(text.split())

        # Truncate if too long
        if len(text) > 500:
            text = text[:500] + "..."

        return text

class ImageDatabase:
    """Manage storage and retrieval of processed images"""

    def __init__(self, data_dir: str = "data/images"):
        self.data_dir = Path(data_dir)
        self.data_dir.mkdir(parents=True, exist_ok=True)

        self.db_file = self.data_dir / "image_database.json"
        self.images = self._load_database()

    def _load_database(self) -> Dict[str, RadiologyImage]:
        """Load image database from JSON file"""
        if not self.db_file.exists():
            return {}

        try:
            with open(self.db_file, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return {img_id: RadiologyImage.from_dict(img_data)
                   for img_id, img_data in data.items()}
        except Exception as e:
            logging.error(f"Error loading image database: {e}")
            return {}

    def _save_database(self):
        """Save image database to JSON file"""
        try:
            data = {img_id: img.to_dict() for img_id, img in self.images.items()}
            with open(self.db_file, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
        except Exception as e:
            logging.error(f"Error saving image database: {e}")

    def add_image(self, image: RadiologyImage):
        """Add image to database"""
        self.images[image.image_id] = image
        self._save_database()

    def add_images(self, images: List[RadiologyImage]):
        """Add multiple images to database"""
        for image in images:
            self.images[image.image_id] = image
        self._save_database()
        logging.info(f"Added {len(images)} images to database")

    def search_images(self, query: str = "", modality: str = "", body_part: str = "",
                     pathology: str = "", limit: int = 20) -> List[RadiologyImage]:
        """Search images by various criteria"""
        results = []

        for image in self.images.values():
            match = True

            # Text search in filename, tags, and extracted text
            if query:
                searchable_text = f"{image.file_path} {' '.join(image.tags)} {image.extracted_text}".lower()
                if query.lower() not in searchable_text:
                    match = False

            # Filter by modality
            if modality and image.modality.lower() != modality.lower():
                match = False

            # Filter by body part
            if body_part and image.body_part.lower() != body_part.lower():
                match = False

            # Filter by pathology
            if pathology and pathology.lower() not in [p.lower() for p in image.pathology]:
                match = False

            if match:
                results.append(image)

        # Sort by relevance (basic implementation)
        if query:
            results.sort(key=lambda x: query.lower() in os.path.basename(x.file_path).lower(), reverse=True)

        return results[:limit]

    def get_stats(self) -> Dict[str, Any]:
        """Get database statistics"""
        total_images = len(self.images)

        # Count by modality
        modality_counts = {}
        body_part_counts = {}

        for image in self.images.values():
            modality_counts[image.modality] = modality_counts.get(image.modality, 0) + 1
            body_part_counts[image.body_part] = body_part_counts.get(image.body_part, 0) + 1

        return {
            "total_images": total_images,
            "modalities": modality_counts,
            "body_parts": body_part_counts,
            "with_embeddings": sum(1 for img in self.images.values() if img.visual_embedding),
            "extracted_from_docs": sum(1 for img in self.images.values() if img.source_document)
        }

class RadiologyImageManager:
    """Main interface for radiology image management"""

    def __init__(self, data_dir: str = "data/images"):
        self.processor = ImageProcessor(data_dir)
        self.extractor = DocumentImageExtractor(os.path.join(data_dir, "extracted"))
        self.database = ImageDatabase(data_dir)

    def scan_directory(self, directory: str, recursive: bool = True) -> int:
        """Scan directory for images and add to database"""
        image_extensions = {'.jpg', '.jpeg', '.png', '.bmp', '.tiff', '.tif', '.dcm'}
        found_images = 0

        try:
            scan_path = Path(directory)
            if not scan_path.exists():
                logging.error(f"Directory does not exist: {directory}")
                return 0

            # Get image files
            if recursive:
                image_files = []
                for ext in image_extensions:
                    image_files.extend(scan_path.rglob(f"*{ext}"))
                    image_files.extend(scan_path.rglob(f"*{ext.upper()}"))
            else:
                image_files = []
                for ext in image_extensions:
                    image_files.extend(scan_path.glob(f"*{ext}"))
                    image_files.extend(scan_path.glob(f"*{ext.upper()}"))

            # Process each image
            new_images = []
            for image_file in image_files:
                # Check if already processed
                with open(str(image_file), 'rb') as f:
                    file_hash = hashlib.md5(f.read()).hexdigest()

                if file_hash not in self.database.images:
                    rad_image = self.processor.process_image(str(image_file))
                    if rad_image:
                        new_images.append(rad_image)
                        found_images += 1

            # Add to database
            if new_images:
                self.database.add_images(new_images)

            logging.info(f"Processed {found_images} new images from {directory}")

        except Exception as e:
            logging.error(f"Error scanning directory {directory}: {e}")

        return found_images

    def extract_from_documents(self, directory: str, doc_types: List[str] = None) -> int:
        """Extract images from documents in directory"""
        if doc_types is None:
            doc_types = ['pdf', 'ppt', 'pptx']

        extracted_count = 0

        try:
            scan_path = Path(directory)
            if not scan_path.exists():
                logging.error(f"Directory does not exist: {directory}")
                return 0

            # Find documents
            for doc_type in doc_types:
                doc_files = list(scan_path.rglob(f"*.{doc_type}"))

                for doc_file in doc_files:
                    try:
                        if doc_type == 'pdf':
                            images = self.extractor.extract_from_pdf(str(doc_file))
                        elif doc_type in ['ppt', 'pptx']:
                            images = self.extractor.extract_from_ppt(str(doc_file))
                        else:
                            continue

                        if images:
                            self.database.add_images(images)
                            extracted_count += len(images)

                    except Exception as e:
                        logging.error(f"Error processing document {doc_file}: {e}")
                        continue

            logging.info(f"Extracted {extracted_count} images from documents")

        except Exception as e:
            logging.error(f"Error extracting from documents: {e}")

        return extracted_count

    def search_images(self, query: str = "", **kwargs) -> List[RadiologyImage]:
        """Search images with various filters"""
        return self.database.search_images(query=query, **kwargs)

    def get_image_by_id(self, image_id: str) -> Optional[RadiologyImage]:
        """Get specific image by ID"""
        return self.database.images.get(image_id)

    def update_image_tags(self, image_id: str, tags: List[str], pathology: List[str] = None):
        """Update image tags and pathology"""
        if image_id in self.database.images:
            self.database.images[image_id].tags = tags
            if pathology:
                self.database.images[image_id].pathology = pathology
            self.database._save_database()

    def get_stats(self) -> Dict[str, Any]:
        """Get comprehensive statistics"""
        return self.database.get_stats()